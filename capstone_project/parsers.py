# parsers.py
import os
import pdfplumber
import openpyxl
from pptx import Presentation
from docx import Document
from config import IMAGE_OUTPUT_DIR
from utils import (
    classify_pdf_document_type,
    analyze_image_with_vlm, analyze_image_with_vlm_fusion, restore_and_analyze_page, get_synthesized_body_text, get_synthesized_body_text2
)

# from docling.document_converter import DocumentConverter
from markitdown import MarkItDown

class AdvancedDocumentParser:
    def __init__(self, file_path, file_hash):
        self.file_path = file_path
        self.file_hash = file_hash
        self.ext = os.path.splitext(file_path)[1].lower()
        self.file_image_dir = os.path.join(IMAGE_OUTPUT_DIR, file_hash)
        os.makedirs(self.file_image_dir, exist_ok=True)
        self.md = MarkItDown()
        
    def parse(self):
        """확장자 매핑 제어 스위치"""
        if self.ext == '.pdf':
            # 💡 유틸 함수에 파일 경로를 던져 가볍게 타입을 알아냅니다.
            doc_type = classify_pdf_document_type(self.file_path)
            if doc_type == "SLIDE_TYPE":
                return self._parse_pdf_as_slide()
            else:
                return self._parse_pdf()
            return self._parse_pdf()
            # return self._parse_pdf_docling()
        elif self.ext in ['.pptx', '.ppt']:
            return self._parse_ppt()
        elif self.ext in ['.docx', '.doc']:
            return self._parse_word()
        elif self.ext in ['.xlsx', '.xls']:
            return self._parse_excel()
        else:
            return self._parse_generic_text()

    #region docling 주석
    # def _get_docling_converter(self) -> DocumentConverter:
    #     """Docling 컨버터 싱글톤 인스턴스 반환"""
    #     global _docling_converter
    #     if _docling_converter is None:
    #         logger.info("Docling AI 모델 및 컨버터를 최초 초기화합니다...")
    #         _docling_converter = DocumentConverter()
    #     return _docling_converter


    # def _get_markitdown_converter(self) -> MarkItDown:
    #     """MarkItDown 컨버터 싱글톤 인스턴스 반환"""
    #     global _markitdown_converter
    #     if _markitdown_converter is None:
    #         logger.info("MarkItDown 컨버터를 최초 초기화합니다...")
    #         _markitdown_converter = MarkItDown()
    #     return _markitdown_converter


    # def _parse_pdf_docling(self) -> str:
    #     """
    #     [내부 함수] Docling을 사용하여 PDF를 고품질 마크다운으로 변환합니다.
    #     """
    #     try:            
    #         converter = self._get_docling_converter()
    #         result = converter.convert(self.file_path)
    #         return result.document.export_to_markdown()
    #     except Exception as e:
    #         logger.error(f"Docling PDF 파싱 실패 ({self.file_path}): {str(e)}")
    #         raise e
    #endregion docling 주석
    
    def _parse_pdf(self):
        import re  # 💡 정규식 매칭을 위한 임포트
        from utils import classify_pdf_document_type  # 💡 앞서 유틸로 뺀 문서 타입 판별기

        parsed_chunks = []
        
        # 1단계: 전체 문서가 PPT(슬라이드) 구조인지 사전 검증
        doc_type = classify_pdf_document_type(self.file_path)
        is_slide_doc = (doc_type == "SLIDE_TYPE")
        
        if is_slide_doc:
            print(f"[*] [강력 권장] {os.path.basename(self.file_path)} 파일은 PPT/슬라이드 스타일로 판정되어 전용 알고리즘을 적용합니다.")

        with pdfplumber.open(self.file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                print(f"[*] [페이지 추출 시작] {os.path.basename(self.file_path)} - {page_num} / {len(pdf.pages)} 페이지 처리 중...")
                native_text = page.extract_text() or ""
                total_physical_glyphs = len(page.chars)
                page_area = page.width * page.height
                
                # 이미지 추출 및 통이미지 판별
                valid_images = []
                is_full_page_image = False
                
                for img in page.images:
                    try:
                        width, height = img.get("width", 0), img.get("height", 0)
                        if width < 40 or height < 40: continue
                        if (width / height) > 12 or (height / width) > 12: continue
                        
                        img_area_ratio = (width * height) / page_area
                        if img_area_ratio > 0.85:
                            is_full_page_image = True
                        valid_images.append(img)
                    except:
                        continue
                
                if total_physical_glyphs < 10 and len(native_text.strip()) == 0:
                    is_full_page_image = True

                fusion_content = ""
                page_saved_images = []  # 현재 페이지에서 물리 디스크에 저장된 이미지 경로를 추적하는 바구니
                section_title = f"Page {page_num}"  # 기본 섹션 타이틀 설정

                # =========================================================
                # [신규 분기] PPT/슬라이드 전용 레이어 (통짜 이미지 렌더링 + 하이브리드 텍스트 융합)
                # =========================================================
                if is_slide_doc:
                    print(f"[*] [Slide Mode] {os.path.basename(self.file_path)} {page_num}p 전체 렌더링 및 VLM 문맥 추출 진행...")
                    full_page_path = os.path.join(self.file_image_dir, f"page_{page_num}_slide_view.png")
                    page.to_image(resolution=150).save(full_page_path, format="PNG")
                    page_saved_images.append(full_page_path)
                    
                    # VLM 경유 전체 레이아웃 및 의미론적 요약 정보 획득
                    vlm_summary = analyze_image_with_vlm_fusion(full_page_path)
                    
                    # 컨텐츠 조립 (시각 레이아웃 요약 + 원본 문자열 백업)
                    fusion_content = f"## [슬라이드 장표 시각 레이아웃 분석 - {page_num}페이지]\n{vlm_summary.strip()}\n\n"
                    if native_text.strip():
                        fusion_content += f"--- [슬라이드 내 원본 텍스트 데이터 (색인 검색 및 무결성 검증용)] ---\n{native_text.strip()}\n"
                    
                    # 슬라이드 첫 줄 문맥을 읽어 섹션 타이틀 메타데이터 고도화
                    lines = [l.strip() for l in native_text.split('\n') if l.strip()]
                    if lines and len(lines[0]) < 40:
                        section_title = f"Page {page_num}: {lines[0]}"

                # =========================================================
                # [분기 0] 일반 문서 내 통이미지 / 스캔본 전용 레이어
                # =========================================================
                elif is_full_page_image:
                    print(f"[*] [통이미지 단일 처리] {os.path.basename(self.file_path)} {page_num}p")
                    full_page_path = os.path.join(self.file_image_dir, f"page_{page_num}_pure_image.png")
                    page.to_image(resolution=150).save(full_page_path, format="PNG")
                    
                    page_saved_images.append(full_page_path)
                    
                    vlm_ocr_result = analyze_image_with_vlm_fusion(full_page_path)
                    fusion_content = f"## [문서 본문(전체 이미지 OCR) - 페이지 {page_num}]\n{vlm_ocr_result}\n"

                # =========================================================
                # [분기 1] 일반/혼합 문서 레이어 (Word, 한글, 보고서형 PDF 등)
                # =========================================================
                else:
                    filtered_images = []
                    for img in valid_images:
                        width, height = img.get("width", 0), img.get("height", 0)
                        if ((width * height) / page_area) > 0.80 and len(native_text.strip()) > 100:
                            continue
                        filtered_images.append(img)
                        
                    filtered_images.sort(key=lambda x: x["top"])
                    is_text_corrupted = self._is_text_corrupted(page, native_text)
                    
                    # ---------------------------------------------------------
                    # CASE A: 텍스트 레이어 오염 -> 기존 구조 활용 전체 복구 후 매핑
                    # ---------------------------------------------------------
                    if is_text_corrupted:
                        print(f"[*] [문맥 인라인 복구] {os.path.basename(self.file_path)} {page_num}p 처리 중...")
                        
                        full_page_path = os.path.join(self.file_image_dir, f"page_{page_num}_layout_repair.png")
                        page.to_image(resolution=150).save(full_page_path, format="PNG")
                        page_saved_images.append(full_page_path)
                        
                        reconstructed_text = get_synthesized_body_text2(full_page_path, native_text, num_images=len(filtered_images))
                        
                        for idx, img in enumerate(filtered_images):
                            img_x0, img_top = max(0, img["x0"]), max(0, img["top"])
                            img_x1, img_bottom = min(page.width, img["x1"]), min(page.height, img["bottom"])
                            
                            if img_x0 < img_x1 and img_top < img_bottom:
                                bbox = (img_x0, img_top, img_x1, img_bottom)
                                cropped = page.crop(bbox).to_image(resolution=150)
                                img_path = os.path.join(self.file_image_dir, f"page_{page_num}_img_{idx}.png")
                                cropped.save(img_path, format="PNG")
                                
                                page_saved_images.append(img_path)
                                img_fusion_result = analyze_image_with_vlm_fusion(img_path)
                                
                                tag = f"[시각자료 #{idx}]"
                                # 💡 전략 2를 반영한 시작/종료 격리 태그 규격화
                                replacement = (
                                    f"\n\n--- [본문 내 시각자료 #{idx} 융합 데이터 시작] ---\n"
                                    f"[시각 자료 인공지능 분석]:\n{img_fusion_result.strip()}\n"
                                    f"--- [본문 내 시각자료 #{idx} 융합 데이터 종료] ---\n\n"
                                )
                                
                                if tag in reconstructed_text:
                                    reconstructed_text = reconstructed_text.replace(tag, replacement)
                                else:
                                    reconstructed_text += f"\n\n{replacement}"
                        
                        fusion_content = f"## [문서 본문 - 페이지 {page_num}]\n{reconstructed_text}\n"

                    # ---------------------------------------------------------
                    # CASE B: 정상 문서 -> 물리적 좌표 기준 순차 조립 (샌드위치 기법 고도화)
                    # ---------------------------------------------------------
                    else:
                        fusion_content = f"## [문서 본문 - 페이지 {page_num}]\n"
                        current_top = 0
                        
                        for idx, img in enumerate(filtered_images):
                            top_coord = max(0, current_top)
                            bottom_coord = min(page.height, img['top'])
                            
                            if top_coord < bottom_coord:
                                slice_bbox = (0, top_coord, page.width, bottom_coord)
                                text_slice = page.crop(slice_bbox).extract_text() or ""
                                if text_slice.strip():
                                    fusion_content += text_slice.strip() + "\n\n"
                            
                            img_x0, img_top = max(0, img["x0"]), max(0, img["top"])
                            img_x1, img_bottom = min(page.width, img["x1"]), min(page.height, img["bottom"])
                            
                            if img_x0 < img_x1 and img_top < img_bottom:
                                bbox = (img_x0, img_top, img_x1, img_bottom)
                                cropped = page.crop(bbox).to_image(resolution=150)
                                img_path = os.path.join(self.file_image_dir, f"page_{page_num}_img_{idx}.png")
                                cropped.save(img_path, format="PNG")
                                
                                page_saved_images.append(img_path)
                                img_fusion_result = analyze_image_with_vlm_fusion(img_path)
                                
                                # 💡 [전략 1 구현] 이미지 하단 캡션 탐색 및 경계선 확장 로직
                                caption_text = ""
                                visual_end_y = img_bottom  # 기본 종료 마커는 이미지 바닥
                                
                                # 이미지 하단 약 30pt 영역 슬라이싱 스캔
                                caption_zone = (0, img_bottom, page.width, min(page.height, img_bottom + 30))
                                potential_caption_text = page.crop(caption_zone).extract_text() or ""
                                
                                # 캡션 정규식 패턴 검사
                                caption_pattern = r'^(그림|표|Fig|Table|※|주\d*)\s*[\d\-\.]+'
                                if re.search(caption_pattern, potential_caption_text.strip()) or any(kw in potential_caption_text for kw in ["출처:", "참조:"]):
                                    caption_text = potential_caption_text.strip()
                                    caption_words = page.crop(caption_zone).extract_words()
                                    if caption_words:
                                        # 감지된 캡션 글자 중 가장 아래에 위치한 y좌표로 경계선 최종 확장
                                        visual_end_y = max(w['bottom'] for w in caption_words)
                                        print(f"[+ Caption Detected] 페이지 {page_num} / 이미지 #{idx} 하단 캡션 인지 성공 -> 경계선 확장 보정")

                                # 💡 [전략 2 구현] 격리 태그 기반 마크다운 덤프
                                fusion_content += f"--- [본문 내 시각자료 #{idx} 융합 데이터 시작] ---\n"
                                if caption_text:
                                    fusion_content += f"[문서 내 공식 캡션]: {caption_text}\n"
                                fusion_content += f"[시각 자료 인공지능 분석]:\n{img_fusion_result.strip()}\n"
                                fusion_content += f"--- [본문 내 시각자료 #{idx} 융합 데이터 종료] ---\n\n"
                                
                                # 💡 포인터를 이미지 밑이 아닌, '캡션 문자 끝'으로 갱신하여 본문 겹침 추출 완전 방어
                                current_top = max(current_top, visual_end_y)
                            else:
                                current_top = max(current_top, img['bottom'])
                        
                        final_top = max(0, current_top)
                        if final_top < page.height:
                            slice_bbox = (0, final_top, page.width, page.height)
                            text_slice = page.crop(slice_bbox).extract_text() or ""
                            if text_slice.strip():
                                fusion_content += text_slice.strip() + "\n"

                # 데이터 스키마 변동 없이 유연하게 적재 추가
                parsed_chunks.append({
                    "content": fusion_content,
                    "page_number": page_num,
                    "section_title": section_title,  # 고도화된 타이틀 매핑
                    "images": page_saved_images
                })
        return parsed_chunks

    def _parse_pdf_as_slide(self):
        """
        [PPT/슬라이드 전용 파싱 알고리즘]
        개별 요소를 쪼개서 크롭하지 않고, 슬라이드 한 장 전체를 고해상도 이미지로 구운 뒤
        VLM의 시각적 문맥 분석 결과와 PDF 내 원본 문자열(Raw Text)을 하이브리드로 결합합니다.
        """
        print(f"[*] 슬라이드형(SLIDE_TYPE) 알고리즘으로 파싱을 시작합니다: {self.file_path}")
        raw_pages = []
        
        # 1. 페이지 통짜 렌더링 이미지를 보관할 디렉토리 검증 및 생성
        if not os.path.exists(self.file_image_dir):
            os.makedirs(self.file_image_dir)
            
        with pdfplumber.open(self.file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                print(f"[*] [Slide Mode] {page_num} / {len(pdf.pages)} 페이지 추출 및 분석 가동...")
                
                # 2. [시각 레이어 1단계] 페이지 전체를 150 DPI PNG 이미지로 통째로 렌더링
                # 💡 resolution=150은 작은 글자 식별력과 VLM 토큰 비용 간의 가장 이상적인 균형점입니다.
                img_path = os.path.join(self.file_image_dir, f"full_page_{page_num}.png")
                try:
                    page.to_image(resolution=150).save(img_path, format="PNG")
                except Exception as img_err:
                    print(f"[-] [Render Error] {page_num}페이지 이미지 변환 실패 (우회 진행): {img_err}")
                    continue

                # 3. [시각 레이어 2단계] VLM(Vision-Language Model)을 경유하여 공간적 배치 및 데이터 흐름 해석
                # 💡 기존 시스템에 구현되어 있는 'analyze_image_with_vlm_fusion' 공용 함수를 활용합니다.
                try:
                    vlm_summary = analyze_image_with_vlm_fusion(img_path)
                except Exception as vlm_err:
                    print(f"[-] [VLM Gateway Error] {page_num}페이지 VLM 분석 실패: {vlm_err}")
                    vlm_summary = "시각 데이터 요약 분석 실패 (LiteLLM Gateway 혹은 AI 인프라 점검 필요)"

                # 4. [텍스트 레이어] 검색 무결성 확보를 위한 순수 Raw 텍스트 추출
                # 💡 VLM이 놓치거나 환각(Hallucination)할 수 있는 미세 숫자, 영어 코드, 품번, 고유명사를 
                # 백업 데이터로 확실히 바인딩하여 하이브리드 검색(FTS) 시 검색 유실율을 0%로 만듭니다.
                raw_text = page.extract_text() or ""

                # 5. 하이브리드 융합 컨텐츠(Fusion Content) 조립
                fusion_content = f"## [슬라이드 분석 - {page_num}페이지]\n"
                fusion_content += f"{vlm_summary.strip()}\n\n"
                
                if raw_text.strip():
                    fusion_content += f"--- [슬라이드 내 원본 텍스트 데이터 (색인 검색 및 무결성 검증용)] ---\n"
                    fusion_content += f"{raw_text.strip()}\n"

                # 6. pipeline.py 및 자식 청커(build_parent_child_packets) 스키마 동기화
                # 슬라이드 첫 줄 텍스트를 감지하여 유의미한 섹션 타이틀(section_title)을 자동 유추합니다.
                section_title = f"제 {page_num}페이지 슬라이드 요약 분석"
                lines = [l.strip() for l in raw_text.split('\n') if l.strip()]
                if lines:
                    potential_title = lines[0]
                    # 첫 줄 텍스트가 너무 길지 않고 유효하다면 섹션 타이틀 힌트로 채택
                    if len(potential_title) < 40:
                        section_title = f"페이지 {page_num}: {potential_title}"

                # 파이프라인 수집 규격 딕셔너리 적재
                raw_pages.append({
                    "page_number": page_num,
                    "section_title": section_title,
                    "content": fusion_content
                })
                
        return raw_pages

    def _is_text_corrupted(self, page, native_text: str) -> bool:
        """
        텍스트 문자열 뒤에 숨은 '무증상 유실(Silent Drop)'을 
        물리적 글자 객체(Glyph) 수와 실제 추출 텍스트의 미스매치로 잡아냅니다.
        """
        # 최종 개선된 조건: 기존 글리프 대비 텍스트 추출률 + 하이브리드 표 오염 감지 + 난수화/특수문자 패턴
        cleaned_text = native_text.strip()
        total_physical_glyphs = len(page.chars)
        
        # -----------------------------------------------------------------
        # [기존 조건 1] 글리프 수 대비 텍스트 추출률이 떨어지는 경우 (폰트 인코딩 깨짐)
        # -----------------------------------------------------------------
        if total_physical_glyphs > 0 and (len(cleaned_text) / total_physical_glyphs) < 0.75:
            return True
            
        # -----------------------------------------------------------------
        # 🔥 [신규 추가] 하이브리드 표 ghost 데이터 누락 감지 로직
        # -----------------------------------------------------------------
        # page.lines(표의 격자선)와 page.rects(표의 셀 테두리/배경)의 총합을 구합니다.
        table_vector_elements = len(page.lines) + len(page.rects)
        
        # 눈에 보이는 표 격자(선/면)는 많은데(예: 15개 이상), 
        # 추출된 텍스트는 헤더 몇 줄 수준(예: 350자 미만)으로 극히 적다면 데이터 누락으로 판단!
        if table_vector_elements > 15 and len(cleaned_text) < 350:
            print(
                f"[*] [하이브리드 오염 감지] {page.page_number}p - "
                f"표 구조({table_vector_elements}개) 대비 텍스트({len(cleaned_text)}자)가 비정상적으로 적습니다. "
                f"VLM 눈 복구(CASE A)로 강제 전환합니다."
            )
            return True
            
        # -----------------------------------------------------------------
        # [기존 조건 2] 난수화 및 특수문자 도배 깨짐 검사
        # -----------------------------------------------------------------
        import re
        if len(cleaned_text) > 150:
            if not re.search(r'[0-9]', cleaned_text) or not re.search(r'[a-zA-Z]', cleaned_text):
                return True
                
        return False

    def _is_likely_table(self, text):
        """내부 텍스트 기호 분석을 통해 표(Table) 구조 유무를 간이 판별"""
        # 공백 연속도가 높거나 숫자, 기호(|, +, -, ┌) 비율이 높은 경우 VLM 분석을 유지하기 위함
        lines = text.split("\n")
        table_score = 0
        for line in lines:
            if "  " in line or "\t" in line or "|" in line:
                table_score += 1
        return table_score > 2

    def _parse_ppt(self):
        """
        MarkItDown과 python-pptx를 하이브리드로 결합하여 
        슬라이드 단위로 페이지를 정밀하게 쪼개고 내부 이미지를 복구합니다.
        """
        parsed_chunks = []
        prs = Presentation(self.file_path)
        # ... (MarkItDown 로드 부분 생략) ...

        for idx, slide in enumerate(prs.slides, start=1):
            slide_text = []
            slide_images = []  # 💡 슬라이드별 이미지 경로를 담을 바구니
            slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {idx}"
            
            slide_text.append(f"## [PPTX 슬라이드 제목: {slide_title}]")
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip(): slide_text.append(paragraph.text)
                
                # 이미지 감지 시
                if shape.shape_type == 13: 
                    try:
                        image = shape.image
                        img_name = f"slide_{idx}_img_{shape.shape_id}.png"
                        img_path = os.path.join(self.file_image_dir, img_name)
                        
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        
                        # 💡 파일 경로 보존 (웹 가상 경로 또는 절대 경로)
                        slide_images.append(img_path)
                        
                        # Gemma 4 VLM 컨텍스트 결합
                        vlm_desc = analyze_image_with_vlm_fusion(img_path)
                        slide_text.append(f"\n[슬라이드 이미지 개체 분석 요약 (Gemma 4)]\n{vlm_desc}")
                    except Exception:
                        continue
                        
            parsed_chunks.append({
                "content": "\n".join(slide_text),
                "page_number": idx,
                "section_title": slide_title,
                "images": slide_images  # 👈 파서 리턴 구조에 추가
            })
        return parsed_chunks

    def _parse_word(self):
        """
        MarkItDown으로 Word 문서의 표와 형식을 마크다운 스트림으로 뽑아낸 후,
        내부 첨부 이미지를 순차적으로 추출하여 매핑합니다.
        """
        parsed_chunks = []
        
        # 1. MarkItDown을 이용한 마크다운 스트림 변환
        try:
            md_result = self.md.convert(self.file_path)
            markdown_text = md_result.text_content
        except Exception as e:
            raise RuntimeError(f"MarkItDown Word 변환 실패: {e}")

        # 2. Word 내부의 이미지 바이너리를 정밀 수집하기 위해 python-docx 사용
        doc = Document(self.file_path)
        docx_images = []
        img_counter = 1
        
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    img_name = f"docx_img_{img_counter}.png"
                    img_path = os.path.join(self.file_image_dir, img_name)
                    
                    with open(img_path, "wb") as f:
                        f.write(rel.target_part.blob)
                    
                    docx_images.append(img_path)
                    
                    # Gemma 4 VLM 컨텍스트 추출 후 부록 형태로 텍스트 스트림 꼬리에 결합
                    vlm_desc = analyze_image_with_vlm_fusion(img_path)
                    markdown_text += f"\n\n[문서 첨부 이미지_{img_counter} 분석 요약 (Gemma 4)]\n{vlm_desc}"
                    img_counter += 1
                except Exception:
                    continue

        # 3. 약 2,000자 단위로 논리적 페이지 섹션을 가상 분할하여 패킷 포맷 충족
        chunk_size = 2000
        current_section = "Word 본문 레이어"
        lines = markdown_text.split("\n")
        
        for i, start in enumerate(range(0, len(markdown_text), chunk_size), start=1):
            chunk_content = markdown_text[start:start+chunk_size]
            
            # 현재 덩어리 영역에서 제목성 텍스트 트래킹 및 메타데이터 동기화
            for line in lines[int(start/100):int((start+chunk_size)/100)+1]:
                if line.startswith("## ") or line.startswith("### "):
                    current_section = line.replace("#", "").strip()
                    break

            parsed_chunks.append({
                "content": f"## [Word 파트 {i}]\n" + chunk_content,
                "page_number": i,
                "section_title": current_section,
                "images": docx_images if i == 1 else [] # 이미지 리스트는 부모 루트 추적용으로 1개 이상 바인딩
            })
            
        return parsed_chunks

    def _parse_excel(self):
        """
        Excel 문서를 시트(Sheet)별로 마크다운 표로 변환하고,
        시트 내부에 삽입된 이미지(차트, 사진 등)를 추출하여 위치 정보와 함께 VLM으로 융합합니다.
        """
        parsed_chunks = []
        
        # 1. MarkItDown을 이용한 엑셀 전체 데이터 마크다운화 (기본 텍스트/표 복구)
        try:
            md_result = self.md.convert(self.file_path)
            excel_markdown = md_result.text_content
        except Exception as e:
            raise RuntimeError(f"MarkItDown Excel 변환 실패: {e}")

        # 2. openpyxl을 로드하여 시트 구조 및 이미지 객체 정밀 탐색
        wb = openpyxl.load_workbook(self.file_path, data_only=True)
        sheet_names = wb.sheetnames
        
        sheets_content = {}
        current_sheet = sheet_names[0] if sheet_names else "Sheet1"
        sheets_content[current_sheet] = []
        
        # MarkItDown 결과물을 시트별로 쪼개기
        for line in excel_markdown.split("\n"):
            matched_sheet = None
            for name in sheet_names:
                if name in line and (line.startswith("#") or line.strip() == name):
                    matched_sheet = name
                    break
            
            if matched_sheet:
                current_sheet = matched_sheet
                if current_sheet not in sheets_content:
                    sheets_content[current_sheet] = []
            else:
                sheets_content[current_sheet].append(line)

        # 3. 각 시트(Sheet) 순회하며 이미지 추출 및 패킷 조립
        for idx, sheet_name in enumerate(sheet_names, start=1):
            sheet_lines = sheets_content.get(sheet_name, [])
            sheet_text = "\n".join(sheet_lines).strip()
            
            # openpyxl에서 실제 시트 객체 가져오기
            ws = wb[sheet_name]
            excel_saved_images = []
            img_counter = 1
            
            # 💡 openpyxl은 시트 내 이미지들을 `ws._images` 리스트에 보관합니다.
            if hasattr(ws, '_images') and ws._images:
                sheet_text += "\n\n### [시트 내 포함된 시각자료 부록]"
                
                for inline_img in ws._images:
                    try:
                        # 이미지 바이너리 추출
                        img_data = inline_img._data()
                        
                        img_name = f"excel_sheet_{idx}_img_{img_counter}.png"
                        img_path = os.path.join(self.file_image_dir, img_name)
                        
                        # 로컬 폴더에 저장
                        with open(img_path, "wb") as f:
                            f.write(img_data)
                            
                        excel_saved_images.append(img_path)
                        
                        # 💡 [고급 기능] 이미지가 위치한 엑셀 셀 좌표(행/열) 알아내기
                        cell_location = ""
                        if inline_img.anchor:
                            try:
                                # 0-indexed를 1-indexed 실좌표로 보정
                                row_idx = inline_img.anchor.from_.row + 1
                                col_idx = inline_img.anchor.from_.col + 1
                                # 알기 쉽게 알파벳 열 좌표로 바꾼다면 openpyxl.utils.get_column_letter 사용 가능
                                cell_location = f" ({row_idx}행 {col_idx}열 부근)"
                            except:
                                pass
                        
                        # Gemma 4 VLM 구동하여 차트/이미지 설명 획득
                        vlm_desc = analyze_image_with_vlm_fusion(img_path)
                        
                        # 텍스트에 이미지 정보 강제 주입 (자식 청킹 시 검색 가치 극대화)
                        sheet_text += f"\n\n--- [엑셀 내 이미지 #{img_counter}{cell_location} 융합 데이터] ---\n{vlm_desc}\n"
                        img_counter += 1
                        
                    except Exception as e:
                        print(f"[-] 엑셀 이미지 추출 중 오류 발생 스킵: {e}")
                        continue

            # 해당 시트에 텍스트도 없고 이미지도 없다면 패스
            if not sheet_text.strip() and not excel_saved_images:
                continue
                
            parsed_chunks.append({
                "content": f"## [Excel 시트명: {sheet_name}]\n" + sheet_text,
                "page_number": idx,  # 시트 번호를 RAG의 가상 페이지 번호로 작동시킴
                "section_title": f"Spreadsheet: {sheet_name}",
                "images": excel_saved_images  # 👈 추출된 이미지 경로 배열 리턴 패킷에 바인딩
            })
            
        wb.close()
        return parsed_chunks
    
    def _parse_generic_text(self):
        with open(self.file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return [{"content": text, "page_number": 1, "section_title": "Document Content"}]
    